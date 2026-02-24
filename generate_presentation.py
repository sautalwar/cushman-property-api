"""Generate PropTracker API Security Demo PowerPoint for Cushman & Wakefield."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
NAVY        = RGBColor(0x1A, 0x23, 0x7E)  # dark navy background
NAVY_DARK   = RGBColor(0x0D, 0x14, 0x5A)  # deeper navy for cards
NAVY_MID    = RGBColor(0x28, 0x3A, 0x9B)  # mid navy accents
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xBB, 0xC5, 0xE0)
VULN_RED    = RGBColor(0xC6, 0x28, 0x28)  # vulnerability accent
FIXED_GREEN = RGBColor(0x2E, 0x7D, 0x32)  # fixed / success
GH_PURPLE   = RGBColor(0x6E, 0x40, 0xC9)  # GitHub Copilot purple
WARN_AMBER  = RGBColor(0xE6, 0x5C, 0x00)  # warning orange
TEAL        = RGBColor(0x00, 0x7A, 0x8A)  # architecture accent
AZURE_BLUE  = RGBColor(0x00, 0x78, 0xD4)  # Azure branding

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def make_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_background(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color,
             line_color=None, line_pt=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", font_size=14, font_color=WHITE, bold=True,
                     line_color=None, line_pt=1.5):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color
            p.font.bold = bold
            p.font.name = "Segoe UI"
            p.alignment = PP_ALIGN.CENTER
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=18, color=WHITE, bold=False,
                alignment=PP_ALIGN.LEFT, font_name="Segoe UI",
                italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=WHITE, font_name="Segoe UI",
                    space_after=8):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(space_after)
    return txBox


def add_divider(slide, left, top, width, color=NAVY_MID, height_pt=3):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(height_pt))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def section_header(slide, title, subtitle="", accent=VULN_RED):
    """Standard slide top-bar: accent divider + title + optional subtitle."""
    add_divider(slide, Inches(0), Inches(0), SLIDE_WIDTH, accent, height_pt=6)
    add_textbox(slide, Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.9),
                title, font_size=34, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.25), Inches(11.5), Inches(0.5),
                    subtitle, font_size=15, color=LIGHT_GRAY)


# ---------------------------------------------------------------------------
# SLIDE 1 — TITLE
# ---------------------------------------------------------------------------
def slide_01(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY_DARK)

    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), GH_PURPLE)

    add_textbox(slide, Inches(1.2), Inches(1.4), Inches(10.5), Inches(1.6),
                "PropTracker: API Security POC",
                font_size=52, color=WHITE, bold=True)

    add_textbox(slide, Inches(1.2), Inches(2.9), Inches(10.5), Inches(0.7),
                "Demonstrating GitHub Actions + GHAS + Copilot Autofix",
                font_size=24, color=GH_PURPLE)

    add_divider(slide, Inches(1.2), Inches(3.65), Inches(6), GH_PURPLE, height_pt=3)

    add_textbox(slide, Inches(1.2), Inches(3.9), Inches(10), Inches(0.55),
                "Property Management + Contractor Marketplace API  "
                "x  10 OWASP API Top 10 Vulnerabilities",
                font_size=16, color=LIGHT_GRAY)

    add_rounded_rect(slide, Inches(1.2), Inches(5.0), Inches(5.5), Inches(0.55),
                     NAVY_MID,
                     "Presented by: [Your Name]  |  GitHub, Inc.",
                     font_size=14, bold=False)

    add_textbox(slide, Inches(1.2), Inches(5.75), Inches(10), Inches(0.5),
                "Cushman & Wakefield  x  Development Engineering Director",
                font_size=14, color=LIGHT_GRAY)

    for i, (label, clr) in enumerate([
        ("GHAS", VULN_RED),
        ("Copilot Autofix", GH_PURPLE),
        ("GitHub Actions", FIXED_GREEN),
        ("Azure", AZURE_BLUE),
    ]):
        add_rounded_rect(slide, Inches(1.2 + i * 2.9), Inches(6.7),
                         Inches(2.5), Inches(0.45),
                         clr, label, font_size=13)

    add_speaker_notes(slide,
        "Welcome -- today we are walking through a hands-on proof of concept built specifically "
        "for Cushman & Wakefield's engineering team. PropTracker is a realistic property "
        "management API we deliberately seeded with all 10 OWASP API Security Top 10 "
        "vulnerabilities. We will show how GitHub Actions, GHAS CodeQL, and Copilot Autofix "
        "find and fix every one of them -- automatically, in your CI/CD pipeline.")


# ---------------------------------------------------------------------------
# SLIDE 2 — THE PROBLEM
# ---------------------------------------------------------------------------
def slide_02(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY)
    section_header(slide,
                   "API Vulnerabilities Cost Companies Millions",
                   "The OWASP API Security Top 10 explains why -- PropTracker proves it",
                   VULN_RED)

    stats = [
        ("$4.88M",  "Average cost of a data breach (IBM, 2024)"),
        ("83%",     "Of breaches involve application-layer APIs"),
        ("10x",     "Faster to fix in code than in production"),
    ]
    for i, (number, label) in enumerate(stats):
        x = Inches(0.7 + i * 4.15)
        add_rounded_rect(slide, x, Inches(2.1), Inches(3.8), Inches(1.1),
                         VULN_RED, "", font_size=14)
        add_textbox(slide, x + Inches(0.1), Inches(2.15), Inches(3.6), Inches(0.7),
                    number, font_size=40, color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), Inches(2.85), Inches(3.6), Inches(0.4),
                    label, font_size=12, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)

    owasp_items = [
        "OWASP API Security Top 10 (2023) covers the most critical API-specific risks",
        "Broken Object Level Authorization  x  Broken Auth  x  Excessive Data Exposure",
        "Rate Limiting  x  Function-Level Auth  x  Mass Assignment  x  Security Misconfiguration",
        "Injection  x  Improper Asset Management  x  Insufficient Logging & Monitoring",
        "Commercial real estate: high-value target -- tenant PII, lease data, contractor financials",
    ]
    add_textbox(slide, Inches(0.7), Inches(3.45), Inches(11.5), Inches(0.45),
                "WHY THIS MATTERS FOR CUSHMAN & WAKEFIELD",
                font_size=13, color=WARN_AMBER, bold=True)
    add_bullet_list(slide, Inches(0.7), Inches(3.85), Inches(11.8), Inches(2.6),
                    owasp_items, font_size=15, color=WHITE, space_after=6)

    add_speaker_notes(slide,
        "The numbers are not academic -- a single API breach in commercial real estate exposes "
        "tenant PII, lease terms, contractor rates, and financial data. "
        "OWASP's API Top 10 is the industry standard for what attackers target first. "
        "We built PropTracker to demonstrate every one of those attack vectors in a safe, "
        "controlled environment so your team can see -- and fix -- them before they hit production.")


# ---------------------------------------------------------------------------
# SLIDE 3 — WHAT WE BUILT
# ---------------------------------------------------------------------------
def slide_03(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY)
    section_header(slide,
                   "What We Built -- PropTracker",
                   "A realistic property management + contractor marketplace API",
                   FIXED_GREEN)

    arch = [
        ("React\nFrontend",   TEAL),
        ("Express\nREST API", GH_PURPLE),
        ("PostgreSQL\nDatabase", AZURE_BLUE),
    ]
    box_w, box_h = Inches(2.6), Inches(1.1)
    gap = Inches(0.7)
    total = len(arch) * box_w + (len(arch) - 1) * (gap + Inches(0.3))
    x0 = (SLIDE_WIDTH - total) / 2

    for i, (label, clr) in enumerate(arch):
        x = x0 + i * (box_w + gap + Inches(0.3))
        add_rounded_rect(slide, x, Inches(2.3), box_w, box_h, clr, label,
                         font_size=16, bold=True)
        if i < len(arch) - 1:
            add_textbox(slide, x + box_w, Inches(2.65),
                        Inches(0.7), Inches(0.4),
                        "->", font_size=26, color=WHITE,
                        alignment=PP_ALIGN.CENTER)

    entities = [
        ("Properties", [
            "Property listings: address, type, market value",
            "Owner + tenant associations and lease terms",
            "Rent roll, TI allowances, lease expiry dates",
            "Maintenance request tracking and history",
        ], TEAL),
        ("Contractors", [
            "Contractor profiles + trade specialties",
            "Licensing, insurance, and bond records",
            "Rating, review, and repeat-hire tracking",
            "Banking details and payout / ACH records",
        ], GH_PURPLE),
        ("Jobs & Bids", [
            "Work orders created by property managers",
            "Contractor bid submission and comparison",
            "Status workflow: open -> awarded -> complete",
            "Invoice generation and payment processing",
        ], WARN_AMBER),
    ]
    card_w = Inches(3.9)
    card_gap = Inches(0.27)
    cx0 = (SLIDE_WIDTH - (3 * card_w + 2 * card_gap)) / 2

    for i, (title, items, clr) in enumerate(entities):
        x = cx0 + i * (card_w + card_gap)
        add_rounded_rect(slide, x, Inches(3.65), card_w, Inches(0.5),
                         clr, title, font_size=15, bold=True)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.2), card_w, Inches(2.7))
        card.fill.solid()
        card.fill.fore_color.rgb = NAVY_DARK
        card.line.color.rgb = clr
        card.line.width = Pt(1.5)
        tf = card.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = "  * " + item
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.name = "Segoe UI"
            p.space_after = Pt(5)

    add_speaker_notes(slide,
        "PropTracker mimics a real-world commercial real estate platform -- the kind Cushman & "
        "Wakefield would build or maintain internally. It has a React dashboard, an Express REST "
        "API with JWT auth, and a PostgreSQL database with realistic schema. "
        "The three core entities give us rich attack surface: IDOR on properties, mass assignment "
        "on contractor profiles, SQL injection on job search, SSRF on job import, and more. "
        "Every vulnerability is realistic -- not contrived.")


# ---------------------------------------------------------------------------
# SLIDE 4 — THE 10 VULNERABILITIES (OWASP MAP)
# ---------------------------------------------------------------------------
def slide_04(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY)
    section_header(slide,
                   "The 10 Vulnerabilities We Embedded",
                   "Every OWASP API Security Top 10 category mapped to PropTracker endpoints",
                   VULN_RED)

    vulns = [
        ("API1",  "BOLA / IDOR",               "/properties/{id}  -> read any property"),
        ("API2",  "Broken Authentication",      "/auth/login  -> brute-force, no lockout"),
        ("API3",  "Broken Object Prop. Auth",   "PATCH /contractors  -> mass assignment"),
        ("API4",  "Unrestricted Resource",      "GET /jobs  -> no pagination, DoS possible"),
        ("API5",  "Broken Func. Level Auth",    "/admin/users  -> no role check"),
        ("API6",  "Unrestricted Access",        "GET /contractors  -> no rate limit"),
        ("API7",  "Server-Side Request Forge",  "POST /jobs/import  -> SSRF via URL param"),
        ("API8",  "Security Misconfiguration",  "CORS wildcard origin (*) on all routes"),
        ("API9",  "Improper Inventory Mgmt",    "Undocumented /debug route exposed in prod"),
        ("API10", "Unsafe Consumption (SQLi)",  "/properties?search=  -> raw SQL injection"),
    ]

    col_w = Inches(6.1)
    row_h = Inches(0.52)
    left1 = Inches(0.55)
    left2 = Inches(6.8)
    top0  = Inches(1.85)

    for i, (api_num, name, endpoint) in enumerate(vulns):
        row = i % 5
        col = i // 5
        left = left1 if col == 0 else left2
        top  = top0 + row * (row_h + Inches(0.06))

        badge_clr = VULN_RED if col == 0 else WARN_AMBER
        add_rounded_rect(slide, left, top, Inches(0.75), Inches(0.48),
                         badge_clr, api_num, font_size=11, bold=True)
        add_textbox(slide, left + Inches(0.85), top, Inches(2.1), Inches(0.48),
                    name, font_size=13, color=WHITE, bold=True)
        add_textbox(slide, left + Inches(3.0), top, Inches(3.0), Inches(0.48),
                    endpoint, font_size=11, color=LIGHT_GRAY, italic=True)

    add_divider(slide, Inches(0.5), Inches(7.05), Inches(12.3), GH_PURPLE, height_pt=2)
    add_textbox(slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
                "Each vulnerability has: a GitHub Actions test step  x  GHAS CodeQL annotation"
                "  x  Copilot Autofix suggestion",
                font_size=13, color=GH_PURPLE, alignment=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "We did not cherry-pick easy vulnerabilities -- we implemented the full OWASP API Top 10. "
        "Each one maps to a real endpoint in PropTracker with a real attack vector. "
        "The left column covers authorization and auth failures; the right covers resource abuse, "
        "configuration issues, and injection. Every row becomes a failing GitHub Actions check "
        "that Copilot Autofix then proposes a code fix for. "
        "The goal: every red check becomes green before the workshop ends.")


# ---------------------------------------------------------------------------
# SLIDE 5 — VULN-1: BOLA
# ---------------------------------------------------------------------------
def slide_05(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY)
    section_header(slide,
                   "VULN-1: Broken Object Level Authorization (BOLA)",
                   "API1  x  The #1 most exploited API vulnerability worldwide",
                   VULN_RED)

    add_textbox(slide, Inches(0.7), Inches(1.75), Inches(4.5), Inches(0.4),
                "WHAT IS IT?", font_size=13, color=VULN_RED, bold=True)
    add_bullet_list(slide, Inches(0.7), Inches(2.1), Inches(4.5), Inches(1.8), [
        "User A can access User B's objects by guessing the ID",
        "API trusts the ID in the URL -- not the session token",
        "Also called IDOR (Insecure Direct Object Reference)",
        "PropTracker: GET /api/properties/42 returns any",
        "  property regardless of who owns it",
    ], font_size=14, color=WHITE, space_after=5)

    add_textbox(slide, Inches(0.7), Inches(3.9), Inches(4.5), Inches(0.4),
                "WHAT HAPPENS IN PROPTRACKER?", font_size=13, color=VULN_RED, bold=True)
    add_bullet_list(slide, Inches(0.7), Inches(4.25), Inches(4.5), Inches(2.1), [
        "Attacker logs in as any contractor (any account)",
        "Iterates: GET /api/properties/1 ... /api/properties/999",
        "Reads lease terms, tenant names, property valuations",
        "No ownership check -- server returns 200 OK every time",
        "All tenant PII and commercial lease data exposed",
    ], font_size=14, color=WHITE, space_after=5)

    add_textbox(slide, Inches(5.7), Inches(1.75), Inches(7.0), Inches(0.4),
                "GITHUB ACTIONS -- WORKFLOW ANNOTATION",
                font_size=13, color=FIXED_GREEN, bold=True)

    code_lines = [
        ("x  VULN-1-BOLA                          [FAILED]", VULN_RED),
        ("", WHITE),
        ("  -> GET /api/properties/99  (user: contractor-2)", LIGHT_GRAY),
        ("     Expected: 403 Forbidden", LIGHT_GRAY),
        ("     Actual:   200 OK  (owner: property-manager-1)", WARN_AMBER),
        ("", WHITE),
        ("  GHAS CodeQL: CWE-639  Authorization Bypass  [HIGH]", WARN_AMBER),
        ("  Copilot Autofix: ownership check in", GH_PURPLE),
        ("    propertyController.getById()  -- 3 lines", GH_PURPLE),
    ]
    code_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.7), Inches(2.15), Inches(7.0), Inches(3.5))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x2A)
    code_box.line.color.rgb = VULN_RED
    code_box.line.width = Pt(1.5)
    tf = code_box.text_frame
    tf.word_wrap = False
    for j, (line, clr) in enumerate(code_lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = clr
        p.font.name = "Courier New"

    add_textbox(slide, Inches(5.7), Inches(5.85), Inches(7.0), Inches(0.4),
                "COPILOT AUTOFIX SUGGESTION",
                font_size=13, color=GH_PURPLE, bold=True)
    fix_lines = [
        ("  if (property.ownerId !== req.user.id) {", FIXED_GREEN),
        ("    return res.status(403).json({ error: 'Forbidden' });", FIXED_GREEN),
        ("  }", FIXED_GREEN),
    ]
    fix_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.7), Inches(6.25), Inches(7.0), Inches(0.85))
    fix_box.fill.solid()
    fix_box.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x2A)
    fix_box.line.color.rgb = GH_PURPLE
    fix_box.line.width = Pt(1.5)
    tf2 = fix_box.text_frame
    tf2.word_wrap = False
    for j, (line, clr) in enumerate(fix_lines):
        p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = clr
        p.font.name = "Courier New"

    add_speaker_notes(slide,
        "BOLA is consistently the #1 API vulnerability because it is so easy to miss in code "
        "review -- the endpoint works correctly for the happy path, just not when another user's "
        "ID is supplied. In PropTracker we simply omit the ownership check in the property "
        "controller. GitHub Actions runs a scripted attacker scenario, gets a 200 when it expects "
        "a 403, and fails the check. GHAS CodeQL annotates the exact line. Copilot Autofix "
        "generates the three-line ownership guard -- developer clicks Accept and opens a PR.")


# ---------------------------------------------------------------------------
# SLIDE 6 — VULN-2 & 3: Broken Auth + Mass Assignment
# ---------------------------------------------------------------------------
def slide_06(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY)
    section_header(slide,
                   "VULN-2 & 3: Broken Authentication + Mass Assignment",
                   "API2  x  API3  x  Auth failures and unintended field overwrites",
                   VULN_RED)

    panels = [
        ("API2 -- Broken Authentication", [
            "Login endpoint: no rate limit, no account lockout",
            "Brute-force 1000 passwords -> eventual compromise",
            "JWT secret hardcoded as 'secret' in source code",
            "Tokens never expire (expiresIn option not set)",
            "Actions check: POST /auth/login x 100 -> expect 429",
            "CodeQL: CWE-307 Improper Restriction of Attempts",
        ], VULN_RED, "Fix: bcrypt + express-rate-limit + env JWT_SECRET + expiry"),
        ("API3 -- Mass Assignment (Broken Object Property Auth)", [
            "PATCH /api/contractors/:id accepts any JSON body",
            "Attacker sends: { \"role\": \"admin\", \"verified\": true }",
            "No field allowlist -- all properties written to DB",
            "Contractor silently elevates own privileges",
            "Actions check: PATCH with forbidden fields -> expect 400",
            "CodeQL: CWE-915 Improperly Controlled Modification",
        ], WARN_AMBER, "Fix: pick() allowlist -- permit only name, phone, specialty"),
    ]

    panel_w = Inches(5.9)
    gap = Inches(0.5)
    x0 = Inches(0.6)

    for i, (title, items, clr, fix) in enumerate(panels):
        x = x0 + i * (panel_w + gap)
        add_rounded_rect(slide, x, Inches(2.0), panel_w, Inches(0.5),
                         clr, title, font_size=14, bold=True)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.55), panel_w, Inches(3.0))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY_DARK
        box.line.color.rgb = clr
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = "  * " + item
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.name = "Segoe UI"
            p.space_after = Pt(5)
        add_rounded_rect(slide, x, Inches(5.65), panel_w, Inches(0.6),
                         FIXED_GREEN, fix, font_size=12, bold=False)

    add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.55),
                "Copilot Autofix generates: rate-limit middleware snippet + field allowlist "
                "for both vulns -- one PR fixes both",
                font_size=14, color=GH_PURPLE)

    add_speaker_notes(slide,
        "Authentication failures and mass assignment are often overlooked because they feel "
        "like developer convenience choices -- no lockout feels harmless, accepting all fields "
        "feels flexible. In PropTracker, a contractor can brute-force the login and then elevate "
        "their own role to admin with a single PATCH request. Copilot Autofix tackles both in one "
        "PR: adds express-rate-limit to the auth route and wraps the PATCH handler with a lodash "
        "pick() allowlist.")


# ---------------------------------------------------------------------------
# SLIDE 7 — VULN-4 & 5: DoS + Function-Level Auth
# ---------------------------------------------------------------------------
def slide_07(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY)
    section_header(slide,
                   "VULN-4 & 5: Unrestricted Resource + Function-Level Auth",
                   "API4  x  API5  x  Resource exhaustion and unprotected admin routes",
                   WARN_AMBER)

    add_rounded_rect(slide, Inches(0.6), Inches(2.0), Inches(5.9), Inches(0.5),
                     WARN_AMBER,
                     "API4 -- Unrestricted Resource Consumption (DoS / Pagination)",
                     font_size=13, bold=True)
    dos_items = [
        "GET /api/jobs returns ALL rows -- no LIMIT clause",
        "With 100k jobs: response payload is ~80MB JSON",
        "Single request can saturate App Service memory",
        "No max page size, no timeout, no streaming",
        "Attack: 10 parallel GET /api/jobs -> OOM crash",
        "Actions check: response time + payload size threshold",
        "Fix: enforce ?page=&limit= max 100, default 25",
    ]
    box4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.55),
        Inches(5.9), Inches(3.5))
    box4.fill.solid()
    box4.fill.fore_color.rgb = NAVY_DARK
    box4.line.color.rgb = WARN_AMBER
    box4.line.width = Pt(1.5)
    tf4 = box4.text_frame
    tf4.word_wrap = True
    for j, item in enumerate(dos_items):
        p = tf4.paragraphs[0] if j == 0 else tf4.add_paragraph()
        p.text = "  * " + item
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
        p.space_after = Pt(5)

    add_rounded_rect(slide, Inches(6.8), Inches(2.0), Inches(5.9), Inches(0.5),
                     VULN_RED,
                     "API5 -- Broken Function Level Authorization (Admin Routes)",
                     font_size=13, bold=True)
    auth_items = [
        "POST /api/admin/users -- lists and promotes all users",
        "Route exists but no middleware checks req.user.role",
        "Any authenticated user can promote themselves to admin",
        "DELETE /api/admin/users/:id -- any user deletable",
        "Endpoint is undocumented -- never appears in Swagger",
        "Actions check: non-admin JWT -> expect 403 on /admin/*",
        "Fix: adminOnly() middleware on all /admin routes",
    ]
    box5 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.55),
        Inches(5.9), Inches(3.5))
    box5.fill.solid()
    box5.fill.fore_color.rgb = NAVY_DARK
    box5.line.color.rgb = VULN_RED
    box5.line.width = Pt(1.5)
    tf5 = box5.text_frame
    tf5.word_wrap = True
    for j, item in enumerate(auth_items):
        p = tf5.paragraphs[0] if j == 0 else tf5.add_paragraph()
        p.text = "  * " + item
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
        p.space_after = Pt(5)

    add_textbox(slide, Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.55),
                "Both caught by scheduled GitHub Actions scans  x  "
                "Autofix adds pagination defaults and adminOnly() middleware in one PR",
                font_size=14, color=GH_PURPLE)

    add_speaker_notes(slide,
        "Resource exhaustion vulnerabilities are easy to miss because they do not cause auth "
        "failures -- they cause outages. In PropTracker, a single client can fetch all 100k job "
        "records and crash the App Service. The admin route issue is even more dangerous: "
        "the route exists, works perfectly, but was added during debugging and never secured. "
        "GitHub Actions catches both with scripted load and role-bypass tests on every push.")


# ---------------------------------------------------------------------------
# SLIDE 8 — VULN-6 & 7: Rate Limit + CORS
# ---------------------------------------------------------------------------
def slide_08(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY)
    section_header(slide,
                   "VULN-6 & 8: No Rate Limiting + CORS Misconfiguration",
                   "API6  x  API8  x  Abuse enablement and cross-origin data leakage",
                   WARN_AMBER)

    panels = [
        ("API6 -- Unrestricted Access to Sensitive Business Flows", [
            "GET /api/contractors has no rate limit applied",
            "Automated scraper harvests all 5,000 contractor profiles",
            "Competitor intelligence: rates, specialties, client lists",
            "No IP throttle, no user-agent check, no CAPTCHA",
            "300 requests/second sustainable indefinitely",
            "Actions check: 500 req burst -> expect 429 after 100",
        ], WARN_AMBER, "Fix: express-rate-limit 100/min per IP on /api/contractors"),
        ("API8 -- Security Misconfiguration (CORS Wildcard)", [
            "app.use(cors()) -- allows ALL origins, no restriction",
            "Malicious site reads /api/properties with victim cookie",
            "Cross-site credential theft with no user interaction",
            "Affects every endpoint including admin routes",
            "Browsers enforce CORS only when server restricts it",
            "Actions check: preflight from evil.example.com -> 200",
        ], VULN_RED, "Fix: cors({ origin: ['https://app.proptracker.com'] })"),
    ]

    panel_w = Inches(5.9)
    gap = Inches(0.5)
    x0 = Inches(0.6)

    for i, (title, items, clr, fix) in enumerate(panels):
        x = x0 + i * (panel_w + gap)
        add_rounded_rect(slide, x, Inches(2.0), panel_w, Inches(0.5),
                         clr, title, font_size=13, bold=True)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.55), panel_w, Inches(3.1))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY_DARK
        box.line.color.rgb = clr
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = "  * " + item
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.name = "Segoe UI"
            p.space_after = Pt(5)
        add_rounded_rect(slide, x, Inches(5.75), panel_w, Inches(0.55),
                         FIXED_GREEN, fix, font_size=11, bold=False)

    add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.55),
                "Copilot Autofix supplies the exact cors() config and rate-limit middleware -- "
                "no googling required",
                font_size=14, color=GH_PURPLE)

    add_speaker_notes(slide,
        "Rate limiting and CORS feel like infrastructure concerns, but they are application-level "
        "choices that developers make every day. cors() with no arguments is a common copy-paste "
        "mistake from Stack Overflow. No rate limit is the default until someone adds it. "
        "PropTracker ships both mistakes by default. GitHub Actions validates them on every push "
        "and Copilot Autofix provides the exact configuration strings needed to fix them.")


# ---------------------------------------------------------------------------
# SLIDE 9 — VULN-8 (SQL Injection)
# ---------------------------------------------------------------------------
def slide_09(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY)
    section_header(slide,
                   "VULN-10: SQL Injection -- Tenant Data at Risk",
                   "API10  x  Unsafe Consumption  x  Most business-critical vulnerability for CRE",
                   VULN_RED)

    add_textbox(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(0.4),
                "THE ATTACK", font_size=13, color=VULN_RED, bold=True)

    sqli_lines = [
        ("GET /api/properties?search=office", LIGHT_GRAY),
        ("  => SELECT * FROM properties WHERE name LIKE '%office%'", FIXED_GREEN),
        ("", WHITE),
        ("Attacker sends:", LIGHT_GRAY),
        ("  ?search=' OR '1'='1", VULN_RED),
        ("  => ... WHERE name LIKE '%' OR '1'='1%'", WARN_AMBER),
        ("  => Returns EVERY row in properties table", WARN_AMBER),
        ("", WHITE),
        ("UNION exfil:", LIGHT_GRAY),
        ("  ?search=' UNION SELECT username,password,", VULN_RED),
        ("           email,null FROM users --", VULN_RED),
        ("  => Dumps entire users table in JSON response", WARN_AMBER),
    ]
    sqli_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.2),
        Inches(5.9), Inches(4.2))
    sqli_box.fill.solid()
    sqli_box.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x2A)
    sqli_box.line.color.rgb = VULN_RED
    sqli_box.line.width = Pt(1.5)
    tf = sqli_box.text_frame
    tf.word_wrap = False
    for j, (line, clr) in enumerate(sqli_lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = clr
        p.font.name = "Courier New"

    add_textbox(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.4),
                "COMMERCIAL REAL ESTATE IMPACT", font_size=13, color=VULN_RED, bold=True)
    impact_items = [
        "Tenant PII: names, emails, lease expiry dates",
        "Commercial lease terms: rent, TI allowances",
        "Contractor banking details and payout amounts",
        "Competitor intelligence on portfolio valuations",
        "GDPR / CCPA breach notification obligation",
        "Average CRE data breach: $6.2M (Ponemon, 2024)",
    ]
    add_bullet_list(slide, Inches(7.0), Inches(2.2), Inches(5.8), Inches(2.4),
                    impact_items, font_size=14, color=WHITE, space_after=6)

    add_textbox(slide, Inches(7.0), Inches(4.7), Inches(5.8), Inches(0.4),
                "GITHUB ACTIONS + GHAS RESPONSE", font_size=13, color=FIXED_GREEN, bold=True)
    fix_lines = [
        ("v  Actions: sqlmap-lite probe on /properties?search=", FIXED_GREEN),
        ("v  CodeQL: CWE-89 SQL Injection  [CRITICAL]", FIXED_GREEN),
        ("v  Autofix: parameterized query -- 3 lines replaced", FIXED_GREEN),
        ("", WHITE),
        ("  BEFORE: `WHERE name LIKE '%${search}%'`", VULN_RED),
        ("  AFTER:  `WHERE name LIKE $1`, [`%${search}%`]", FIXED_GREEN),
    ]
    fix_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(5.15),
        Inches(5.8), Inches(1.9))
    fix_box.fill.solid()
    fix_box.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x2A)
    fix_box.line.color.rgb = FIXED_GREEN
    fix_box.line.width = Pt(1.5)
    tf2 = fix_box.text_frame
    tf2.word_wrap = False
    for j, (line, clr) in enumerate(fix_lines):
        p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = clr
        p.font.name = "Courier New"

    add_speaker_notes(slide,
        "SQL injection is decades old but still the #1 cause of large-scale data breaches. "
        "For Cushman & Wakefield the impact is not just technical -- it is tenant PII, lease "
        "confidentiality, contractor financials, and regulatory exposure under GDPR and CCPA. "
        "In PropTracker we use raw string interpolation in one search query -- exactly the kind "
        "of thing that slips through in a fast-moving codebase. CodeQL catches it at the AST "
        "level (semantic analysis, not pattern matching), and Copilot Autofix replaces the "
        "interpolation with a parameterized query in under 10 seconds.")


# ---------------------------------------------------------------------------
# SLIDE 10 — VULN-9 & 10: SSRF + Debug Route
# ---------------------------------------------------------------------------
def slide_10(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY)
    section_header(slide,
                   "VULN-7 & 9: SSRF + Improper Inventory Management",
                   "API7  x  API9  x  Server hijacking and undocumented debug routes",
                   VULN_RED)

    panels = [
        ("API7 -- Server-Side Request Forgery (SSRF)", [
            "POST /api/jobs/import accepts { url: '...' }",
            "Server fetches the URL and parses the JSON body",
            "Attack: { url: 'http://169.254.169.254/metadata' }",
            "Azure IMDS returns VM identity + access tokens",
            "Attacker harvests subscription-level credentials",
            "Actions check: SSRF probe to 169.254.x.x -> fail on 200",
            "Fix: URL allowlist + block all RFC-1918 address ranges",
        ], VULN_RED, "Critical in Azure: IMDS leaks managed identity tokens"),
        ("API9 -- Improper Inventory Management (Debug Route)", [
            "GET /api/debug returns env vars + DB connection string",
            "Route added during dev sprint, never removed",
            "No authentication required -- completely open to anyone",
            "Exposes: DATABASE_URL, JWT_SECRET, STRIPE_KEY",
            "All three secrets directly usable for further attacks",
            "Actions check: GET /api/debug -> expect 404, got 200",
            "Fix: remove route; add route inventory audit to CI",
        ], WARN_AMBER, "Often forgotten: old debug routes left in production code"),
    ]

    panel_w = Inches(5.9)
    gap = Inches(0.5)
    x0 = Inches(0.6)

    for i, (title, items, clr, note) in enumerate(panels):
        x = x0 + i * (panel_w + gap)
        add_rounded_rect(slide, x, Inches(2.0), panel_w, Inches(0.5),
                         clr, title, font_size=13, bold=True)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.55), panel_w, Inches(3.2))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY_DARK
        box.line.color.rgb = clr
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = "  * " + item
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.name = "Segoe UI"
            p.space_after = Pt(5)
        add_rounded_rect(slide, x, Inches(5.85), panel_w, Inches(0.55),
                         WARN_AMBER, note, font_size=11, bold=False,
                         font_color=RGBColor(0x0A, 0x0E, 0x2A))

    add_textbox(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.5),
                "Copilot Autofix: adds SSRF URL blocklist helper + flags debug route for "
                "deletion -- PR description includes CVE reference",
                font_size=14, color=GH_PURPLE)

    add_speaker_notes(slide,
        "SSRF is particularly dangerous in Azure because the Instance Metadata Service endpoint "
        "returns managed identity access tokens -- giving an attacker Azure subscription-level "
        "access from a single malformed API request. The debug route is a code hygiene issue "
        "that becomes a critical vulnerability the moment it hits production. "
        "Both are found by GitHub Actions behavioral tests (not just static analysis), making "
        "them perfect complements to CodeQL's static coverage.")


# ---------------------------------------------------------------------------
# SLIDE 11 — GITHUB ACTIONS: HOW IT WORKS
# ---------------------------------------------------------------------------
def slide_11(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY)
    section_header(slide,
                   "GitHub Actions -- Automated Security Scanning",
                   "Every push and nightly schedule triggers the full vulnerability test suite",
                   FIXED_GREEN)

    steps = [
        ("1  Build & Deploy to Staging",   FIXED_GREEN,
         "docker build -> push to ACR -> deploy to App Service staging slot"),
        ("2  GHAS CodeQL Analysis",         GH_PURPLE,
         "CodeQL scans all 10 CWEs: CWE-89, CWE-639, CWE-307 + more"),
        ("3  OWASP API Security Tests",     VULN_RED,
         "10 scripted attack scenarios -> HTTP status code assertions"),
        ("4  Annotate PR / Step Summary",   WARN_AMBER,
         "Failing vulns surface as PR annotations + markdown job summary"),
        ("5  Copilot Autofix Trigger",      GH_PURPLE,
         "Each CodeQL alert triggers an AI-generated Autofix suggestion"),
        ("6  Security Gate",                FIXED_GREEN,
         "Branch protection blocks merge until all 10 checks pass"),
    ]

    step_h   = Inches(0.6)
    step_gap = Inches(0.07)
    x_left   = Inches(0.6)
    y_start  = Inches(2.0)

    for i, (label, clr, detail) in enumerate(steps):
        y = y_start + i * (step_h + step_gap)
        add_rounded_rect(slide, x_left, y, Inches(0.55), step_h,
                         clr, str(i + 1), font_size=18, bold=True)
        add_textbox(slide, x_left + Inches(0.7), y + Inches(0.1),
                    Inches(4.5), step_h,
                    label, font_size=15, color=WHITE, bold=True)
        add_textbox(slide, x_left + Inches(5.4), y + Inches(0.12),
                    Inches(7.0), step_h,
                    detail, font_size=13, color=LIGHT_GRAY)
        if i < len(steps) - 1:
            add_textbox(slide, x_left + Inches(0.15),
                        y + step_h - Inches(0.01), Inches(0.3), Inches(0.12),
                        "|", font_size=10, color=clr)

    add_rounded_rect(slide, Inches(0.6), Inches(6.65), Inches(5.5), Inches(0.5),
                     NAVY_MID,
                     "Schedule: nightly 02:00 UTC + on: push to main/feature/*",
                     font_size=13, bold=False,
                     line_color=FIXED_GREEN, line_pt=1)
    add_rounded_rect(slide, Inches(6.4), Inches(6.65), Inches(6.3), Inches(0.5),
                     NAVY_MID,
                     "Step Summary: markdown report of all 10 vulns with pass/fail status",
                     font_size=13, bold=False,
                     line_color=GH_PURPLE, line_pt=1)

    add_speaker_notes(slide,
        "The GitHub Actions workflow is the engine behind the entire demo. It runs on every "
        "push and nightly, so you are never more than 24 hours from knowing your security posture. "
        "The pipeline deploys to a staging slot first, runs CodeQL analysis, then executes our "
        "custom OWASP API attack scripts against the live staging API. Failures annotate the PR "
        "directly -- no separate dashboard to check. The security gate means nothing ships until "
        "all 10 vulnerability checks pass.")


# ---------------------------------------------------------------------------
# SLIDE 12 — COPILOT AUTOFIX
# ---------------------------------------------------------------------------
def slide_12(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY)
    section_header(slide,
                   "Copilot Autofix -- AI Suggests the Fix",
                   "From failed annotation to merged PR fix in under 60 seconds",
                   GH_PURPLE)

    flow_steps = [
        ("1", "CodeQL detects CWE-89 in propertyController.js",   VULN_RED),
        ("2", "GitHub creates Security Alert with file + line ref", WARN_AMBER),
        ("3", "Copilot Autofix analyzes AST context + docs",        GH_PURPLE),
        ("4", "AI generates parameterized query replacement",        GH_PURPLE),
        ("5", "Developer sees inline suggestion in PR review",       FIXED_GREEN),
        ("6", "One-click Accept -> commit -> all checks pass",       FIXED_GREEN),
    ]

    for i, (num, text, clr) in enumerate(flow_steps):
        y = Inches(2.05) + i * Inches(0.72)
        add_rounded_rect(slide, Inches(0.6), y, Inches(0.5), Inches(0.55),
                         clr, num, font_size=16, bold=True)
        add_textbox(slide, Inches(1.25), y + Inches(0.07),
                    Inches(5.3), Inches(0.55),
                    text, font_size=14, color=WHITE)

    add_textbox(slide, Inches(7.2), Inches(1.8), Inches(5.8), Inches(0.4),
                "COPILOT AUTOFIX DIFF -- propertyController.js",
                font_size=12, color=GH_PURPLE, bold=True)

    diff_lines = [
        ("- const sql = `SELECT * FROM properties", VULN_RED),
        ("-   WHERE name LIKE '%${req.query.search}%'`;", VULN_RED),
        ("- const result = await db.query(sql);", VULN_RED),
        ("", WHITE),
        ("+ const sql = 'SELECT * FROM properties", FIXED_GREEN),
        ("+   WHERE name LIKE $1';", FIXED_GREEN),
        ("+ const result = await db.query(sql,", FIXED_GREEN),
        ("+   [`%${req.query.search}%`]);", FIXED_GREEN),
    ]
    diff_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(2.25),
        Inches(5.8), Inches(2.5))
    diff_box.fill.solid()
    diff_box.fill.fore_color.rgb = RGBColor(0x04, 0x07, 0x1A)
    diff_box.line.color.rgb = GH_PURPLE
    diff_box.line.width = Pt(1.5)
    tf = diff_box.text_frame
    tf.word_wrap = False
    for j, (line, clr) in enumerate(diff_lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = clr
        p.font.name = "Courier New"

    stat_items = [
        ("7x",    "Faster remediation vs. manual fix (GitHub data)"),
        ("100%",  "PropTracker vulns have Autofix suggestions"),
        ("< 60s", "From CodeQL alert to suggested PR"),
    ]
    for i, (num, label) in enumerate(stat_items):
        y = Inches(4.95) + i * Inches(0.72)
        add_rounded_rect(slide, Inches(7.2), y, Inches(1.2), Inches(0.6),
                         GH_PURPLE, num, font_size=18, bold=True)
        add_textbox(slide, Inches(8.55), y + Inches(0.1),
                    Inches(4.4), Inches(0.5),
                    label, font_size=14, color=WHITE)

    add_speaker_notes(slide,
        "Copilot Autofix closes the loop between detection and remediation. A developer does not "
        "need to understand SQL injection theory -- they see a diff that replaces string "
        "interpolation with a parameterized query, with a plain-English explanation of why it "
        "matters. For PropTracker all 10 vulnerabilities have Autofix suggestions, meaning a new "
        "developer could fix the entire codebase in a single afternoon just by reviewing and "
        "accepting suggested PRs -- no security expertise required.")


# ---------------------------------------------------------------------------
# SLIDE 13 — THE VULNERABILITY DASHBOARD
# ---------------------------------------------------------------------------
def slide_13(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY)
    section_header(slide,
                   "The PropTracker Vulnerability Dashboard",
                   "React frontend -- live GitHub Actions status for all 10 vulnerabilities",
                   TEAL)

    add_textbox(slide, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.4),
                "Each card: vulnerability description  x  attack payload  x  "
                "GitHub Actions status  x  Copilot Autofix availability",
                font_size=14, color=LIGHT_GRAY)

    vuln_cards = [
        ("API1",  "BOLA",              "FAIL", VULN_RED),
        ("API2",  "Broken Auth",       "FAIL", VULN_RED),
        ("API3",  "Mass Assignment",   "FAIL", VULN_RED),
        ("API4",  "DoS / Pagination",  "FAIL", VULN_RED),
        ("API5",  "Func Level Auth",   "FAIL", VULN_RED),
        ("API6",  "Rate Limit",        "FAIL", WARN_AMBER),
        ("API7",  "SSRF",              "FAIL", VULN_RED),
        ("API8",  "CORS Wildcard",     "FAIL", WARN_AMBER),
        ("API9",  "Debug Route",       "FAIL", VULN_RED),
        ("API10", "SQL Injection",     "FAIL", VULN_RED),
    ]

    card_w   = Inches(2.3)
    card_h   = Inches(1.4)
    col_gap  = Inches(0.18)
    row_gap  = Inches(0.18)
    cols     = 5
    total_w  = cols * card_w + (cols - 1) * col_gap
    x0       = (SLIDE_WIDTH - total_w) / 2
    y0       = Inches(2.25)

    for i, (api, name, status, clr) in enumerate(vuln_cards):
        col = i % cols
        row = i // cols
        x = x0 + col * (card_w + col_gap)
        y = y0 + row * (card_h + row_gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = NAVY_DARK
        card.line.color.rgb = clr
        card.line.width = Pt(2)

        add_textbox(slide, x + Inches(0.1), y + Inches(0.08),
                    card_w - Inches(0.2), Inches(0.35),
                    api, font_size=11, color=clr, bold=True)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.4),
                    card_w - Inches(0.2), Inches(0.45),
                    name, font_size=13, color=WHITE, bold=True)

        add_rounded_rect(slide,
                         x + Inches(0.1), y + card_h - Inches(0.42),
                         card_w - Inches(0.2), Inches(0.3),
                         clr, "* " + status, font_size=10, bold=True)

    add_divider(slide, Inches(0.6), Inches(5.95), Inches(12.1), FIXED_GREEN, height_pt=2)
    add_textbox(slide, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.5),
                "After Copilot Autofix PRs are merged -> all cards turn green (PASS)  "
                "x  Dashboard polls GitHub Actions API every 30 seconds",
                font_size=14, color=FIXED_GREEN)
    add_textbox(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.5),
                "Live demo: https://proptracker-poc.azurestaticapps.net  "
                "(available during workshop)",
                font_size=13, color=LIGHT_GRAY, italic=True)

    add_speaker_notes(slide,
        "The dashboard is the showpiece of the demo -- a React app that fetches GitHub Actions "
        "workflow run status via the GitHub API and renders a real-time card for each "
        "vulnerability. Before the demo all 10 cards are red. During the workshop we accept "
        "Copilot Autofix PRs one by one and the audience watches the cards flip to green. "
        "It is a visceral, visual representation of the value: detect, fix, verify -- all in "
        "one pipeline, no external security tool required.")


# ---------------------------------------------------------------------------
# SLIDE 14 — AZURE ARCHITECTURE
# ---------------------------------------------------------------------------
def slide_14(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY)
    section_header(slide,
                   "Azure Architecture -- Full Production Deployment",
                   "GitHub -> Actions -> ACR -> App Service -> PostgreSQL + Static Web App",
                   AZURE_BLUE)

    tiers = [
        ("GitHub\nRepository",            GH_PURPLE,  Inches(0.5),  Inches(2.2)),
        ("GitHub\nActions",               FIXED_GREEN, Inches(2.85), Inches(2.2)),
        ("Azure Container\nRegistry",     AZURE_BLUE,  Inches(5.2),  Inches(2.2)),
        ("Azure App\nService (API)",      AZURE_BLUE,  Inches(7.55), Inches(2.2)),
        ("PostgreSQL\nFlexible Server",   TEAL,        Inches(9.9),  Inches(2.2)),
    ]

    box_w = Inches(2.1)
    box_h = Inches(1.2)

    for i, (label, clr, x, y) in enumerate(tiers):
        add_rounded_rect(slide, x, y, box_w, box_h, clr, label,
                         font_size=13, bold=True)
        if i < len(tiers) - 1:
            next_x = tiers[i + 1][2]
            mid_x = x + box_w + Inches(0.05)
            add_textbox(slide, mid_x, y + Inches(0.38),
                        Inches(0.3), Inches(0.45),
                        "->", font_size=18, color=WHITE,
                        alignment=PP_ALIGN.CENTER)

    # Static Web App below
    add_rounded_rect(slide, Inches(5.2), Inches(3.65), Inches(4.45), Inches(0.8),
                     TEAL,
                     "Azure Static Web Apps  x  React Vulnerability Dashboard",
                     font_size=13, bold=True)
    add_textbox(slide, Inches(7.0), Inches(3.4), Inches(0.3), Inches(0.28),
                "^", font_size=16, color=TEAL, alignment=PP_ALIGN.CENTER)

    # Security layer overlays
    add_rounded_rect(slide, Inches(0.5), Inches(4.7), Inches(5.7), Inches(0.6),
                     RGBColor(0x0A, 0x0E, 0x2A),
                     "GHAS: CodeQL  x  Secret Scanning  x  Dependency Review",
                     font_size=13, bold=False,
                     line_color=VULN_RED, line_pt=1.5)
    add_rounded_rect(slide, Inches(6.5), Inches(4.7), Inches(6.3), Inches(0.6),
                     RGBColor(0x0A, 0x0E, 0x2A),
                     "Copilot Autofix  x  PR Security Gate  x  Branch Protection",
                     font_size=13, bold=False,
                     line_color=GH_PURPLE, line_pt=1.5)

    infra_left = [
        "GitHub Actions: ubuntu-latest runners",
        "Docker multi-stage build (node:20-alpine)",
        "ACR: Basic tier -- image retention 7 days",
        "Deployment slots: staging -> production swap",
    ]
    infra_right = [
        "App Service: B2 (2 vCPU, 3.5 GB) -- burstable",
        "PostgreSQL Flexible Server: Burstable B1ms",
        "VNet integration: API -> DB private endpoint",
        "Static Web Apps: Free tier -- GitHub CI/CD built-in",
    ]
    add_bullet_list(slide, Inches(0.5), Inches(5.5), Inches(5.9), Inches(1.5),
                    infra_left, font_size=13, color=LIGHT_GRAY, space_after=4)
    add_bullet_list(slide, Inches(6.9), Inches(5.5), Inches(5.9), Inches(1.5),
                    infra_right, font_size=13, color=LIGHT_GRAY, space_after=4)

    add_speaker_notes(slide,
        "The entire PropTracker stack is deployed on Azure -- the natural home for a Cushman & "
        "Wakefield workload. GitHub Actions handles CI/CD natively; the workflow builds a Docker "
        "image, pushes to Azure Container Registry, and swaps the staging slot to production "
        "after all security gates pass. The PostgreSQL Flexible Server sits behind a private "
        "endpoint -- only the App Service can reach it. This architecture is production-"
        "representative and can be stood up from the included Bicep templates in under 20 minutes.")


# ---------------------------------------------------------------------------
# SLIDE 15 — GETTING STARTED / CTA
# ---------------------------------------------------------------------------
def slide_15(prs):
    slide = blank_slide(prs)
    add_background(slide, NAVY_DARK)

    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), GH_PURPLE)

    add_textbox(slide, Inches(1.0), Inches(0.45), Inches(11), Inches(0.9),
                "Your API Could Have These Same Issues",
                font_size=38, color=WHITE, bold=True)
    add_textbox(slide, Inches(1.0), Inches(1.3), Inches(11), Inches(0.5),
                "PropTracker is your mirror -- let's find out what's in your codebase",
                font_size=18, color=LIGHT_GRAY, italic=True)

    next_steps = [
        ("Free Trial\n(30 days)",
         "Enable GitHub Copilot Enterprise + GHAS on your org today. "
         "No credit card required. Full feature set. We'll help you configure it.",
         GH_PURPLE),
        ("Security Workshop\n(Half-day on-site)",
         "We run PropTracker-style scans against your own API. "
         "You leave with a prioritized vulnerability report and fix PRs ready to merge.",
         VULN_RED),
        ("Architecture Review\n(60-min call)",
         "Map your Azure + ADO + GitHub environment. "
         "Design your GHAS rollout plan with a GitHub Solutions Engineer.",
         AZURE_BLUE),
    ]

    ns_w   = Inches(3.8)
    ns_gap = Inches(0.35)
    ns_x0  = (SLIDE_WIDTH - (3 * ns_w + 2 * ns_gap)) / 2

    for i, (title, body, clr) in enumerate(next_steps):
        x = ns_x0 + i * (ns_w + ns_gap)
        add_rounded_rect(slide, x, Inches(2.1), ns_w, Inches(0.75),
                         clr, title, font_size=14, bold=True)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.9), ns_w, Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY_MID
        box.line.color.rgb = clr
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = body
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"

    chips = [
        ("v  No migration required",      FIXED_GREEN),
        ("v  Works with Azure DevOps",     AZURE_BLUE),
        ("v  SOC 2 / ISO 27001 certified", GH_PURPLE),
        ("v  IP indemnity included",       WARN_AMBER),
    ]
    for i, (label, clr) in enumerate(chips):
        add_rounded_rect(slide, Inches(1.0 + i * 2.9), Inches(4.95),
                         Inches(2.65), Inches(0.48),
                         clr, label, font_size=12, bold=False)

    add_rounded_rect(slide, Inches(1.0), Inches(5.7), Inches(11.2), Inches(0.65),
                     NAVY_MID,
                     "Repo: github.com/[org]/proptracker-apim-poc  "
                     "x  Contact: [your-email]@github.com  "
                     "x  github.com/features/security",
                     font_size=14, bold=False,
                     line_color=GH_PURPLE, line_pt=1)

    add_textbox(slide, Inches(1.0), Inches(6.55), Inches(11.2), Inches(0.7),
                '"Every vulnerability PropTracker has -- your API might have too. '
                "Let's check. Today.\"",
                font_size=17, color=GH_PURPLE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "Close with urgency -- these vulnerabilities are not hypothetical; they are the same "
        "patterns that caused breaches at companies with engineering teams just like Cushman's. "
        "The three paths forward are ordered by commitment: a free trial costs nothing, a workshop "
        "surfaces real issues in your actual codebase, and an architecture review gives you a "
        "concrete roadmap. Ask: 'Which of these three would be most valuable to your team right "
        "now?' -- then stop talking and let them answer.")


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------
def main():
    import os
    prs = make_prs()
    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)
    slide_13(prs)
    slide_14(prs)
    slide_15(prs)

    out_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "PropTracker-API-Security-Demo.pptx"
    )
    prs.save(out_path)


if __name__ == '__main__':
    main()
    print('Generated: PropTracker-API-Security-Demo.pptx')